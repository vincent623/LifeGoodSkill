"""PPT 生成模块"""
import os
from pptx import Presentation
from pptx.util import Inches


def generate_ppt(emf_files, output_pptx):
    """步骤 5: 创建 PPT 并插入 EMF"""
    print(f"💾 [4/5] 正在生成 PPT 文件: {os.path.basename(output_pptx)}...")

    prs = Presentation()
    # 设置 16:9 宽屏
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    success_count = 0

    for emf_path in emf_files:
        if not emf_path:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

        # 居中插入逻辑
        # 假设 EMF 也是宽屏比例，适当留白
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12.33)
        # height 会自动按比例缩放

        try:
            slide.shapes.add_picture(emf_path, left, top, width=width)
            success_count += 1
            print(f"    -> ✅ 已插入: {os.path.basename(emf_path)}")
        except Exception as e:
            print(f"    ⚠️ 无法插入 {emf_path}: {e}")

    prs.save(output_pptx)
    return success_count
